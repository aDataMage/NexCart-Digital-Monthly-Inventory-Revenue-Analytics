"""
PowerPoint executive summary generator for monthly business reports.
Creates professional presentation with LLM-generated insights and data tables.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


class ExecutiveSummaryGenerator:
    def __init__(self):
        """Initialize PowerPoint generator with professional theme."""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Color scheme (professional blue gradient)
        self.primary_color = RGBColor(52, 152, 219)  # Blue
        self.secondary_color = RGBColor(44, 62, 80)  # Dark gray
        self.accent_color = RGBColor(230, 126, 34)   # Orange
        self.text_color = RGBColor(44, 62, 80)
        self.light_gray = RGBColor(236, 240, 241)

    def generate_presentation(self, insights: dict, report_data: dict, output_path: str):
        """
        Generate complete PowerPoint presentation.
        
        Args:
            insights: LLM-generated insights dictionary
            report_data: Aggregated report data
            output_path: Path to save the .pptx file
        """
        try:
            report_month = report_data.get('report_month', 'Monthly Report')
            
            # Slide 1: Title
            self._create_title_slide(report_month)
            
            # Slide 2: Executive Summary
            self._create_executive_summary_slide(insights.get('executive_summary', ''))
            
            # Slide 3: Sales & Channel Insights
            self._create_sales_insights_slide(
                insights.get('sales_insights', ''),
                report_data.get('sales_summary', {})
            )
            
            # Slide 4: Product & Inventory Insights
            self._create_product_insights_slide(
                insights.get('product_insights', ''),
                report_data.get('product_summary', {})
            )
            
            # Slide 5: Customer Insights
            self._create_customer_insights_slide(
                insights.get('customer_insights', ''),
                report_data.get('customer_summary', {})
            )
            
            # Slide 6: Financial Insights
            self._create_financial_insights_slide(
                insights.get('financial_insights', ''),
                report_data.get('payment_summary', {})
            )
            
            # Slide 7: Recommendations
            self._create_recommendations_slide(insights.get('recommendations', []))
            
            # Slide 8: Risk Alerts
            self._create_risk_alerts_slide(insights.get('risk_alerts', []))
            
            # Save presentation
            self.prs.save(output_path)
            logger.info(f"PowerPoint presentation saved to {output_path}")
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint: {e}")
            raise

    def _create_title_slide(self, report_month: str):
        """Create title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add gradient background (simulated with shapes)
        background = slide.shapes.add_shape(
            1,  # Rectangle
            0, 0,
            self.prs.slide_width,
            self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.primary_color
        background.line.color.rgb = self.primary_color
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary Report"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4),
            Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = report_month
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER

    def _create_executive_summary_slide(self, summary_text: str):
        """Create executive summary slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Title and content
        
        # Set title
        title = slide.shapes.title
        title.text = "Executive Summary"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        # Add content
        content_box = slide.placeholders[1]
        tf = content_box.text_frame
        tf.text = summary_text
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = self.text_color
        tf.word_wrap = True

    def _create_sales_insights_slide(self, insights_text: str, sales_data: dict):
        """Create sales insights slide with data table."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Title only
        
        # Title
        title = slide.shapes.title
        title.text = "Sales & Channel Performance"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        # Insights text
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(2)
        )
        tf = text_box.text_frame
        tf.text = insights_text
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = self.text_color
        tf.word_wrap = True
        
        # Key metrics table
        if sales_data:
            self._add_key_metrics_table(slide, sales_data, top=Inches(3.8))

    def _create_product_insights_slide(self, insights_text: str, product_data: dict):
        """Create product insights slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Product & Inventory Performance"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(5)
        )
        tf = text_box.text_frame
        tf.text = insights_text
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = self.text_color
        tf.word_wrap = True

    def _create_customer_insights_slide(self, insights_text: str, customer_data: dict):
        """Create customer insights slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Customer Demographics & Behavior"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(5)
        )
        tf = text_box.text_frame
        tf.text = insights_text
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = self.text_color
        tf.word_wrap = True

    def _create_financial_insights_slide(self, insights_text: str, financial_data: dict):
        """Create financial insights slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Financial & Payment Analysis"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(5)
        )
        tf = text_box.text_frame
        tf.text = insights_text
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = self.text_color
        tf.word_wrap = True

    def _create_recommendations_slide(self, recommendations: list):
        """Create recommendations slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Strategic Recommendations"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        
        content_box = slide.placeholders[1]
        tf = content_box.text_frame
        tf.clear()
        
        for i, rec in enumerate(recommendations[:6], 1):  # Max 6 recommendations
            p = tf.add_paragraph()
            p.text = f"{i}. {rec}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.text_color
            p.level = 0
            p.space_before = Pt(10)

    def _create_risk_alerts_slide(self, risks: list):
        """Create risk alerts slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "Risk Alerts & Action Items"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = self.accent_color
        
        content_box = slide.placeholders[1]
        tf = content_box.text_frame
        tf.clear()
        
        for risk in risks[:6]:  # Max 6 risks
            p = tf.add_paragraph()
            p.text = f"⚠ {risk}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.text_color
            p.level = 0
            p.space_before = Pt(10)

    def _add_key_metrics_table(self, slide, metrics: dict, top: float):
        """Add a simple key metrics table to a slide."""
        rows = 3
        cols = 2
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.5), top,
            Inches(9), Inches(2.5)
        ).table
        
        # Header styling
        for cell in table.rows[0].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.primary_color
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
        
        # Populate table
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        
        table.cell(1, 0).text = "Total Gross Revenue"
        table.cell(1, 1).text = f"${metrics.get('total_gross', 0):,.2f}"
        
        table.cell(2, 0).text = "Total Orders"
        table.cell(2, 1).text = f"{metrics.get('total_orders', 0):,}"
        
        # Data cell styling
        for row_idx in range(1, rows):
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                if col_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.light_gray


if __name__ == "__main__":
    # Test the generator
    logging.basicConfig(level=logging.INFO)
    
    test_insights = {
        'executive_summary': "Business performance in October 2025 showed strong growth with total revenue of $224K, representing a healthy mix across all sales channels.",
        'sales_insights': "Amazon and Shopify channels drove the majority of revenue, with Amazon showing a 15% commission impact. Net margins remain strong at 91%.",
        'product_insights': "Top 20% of products generated 78% of revenue, indicating a healthy Pareto distribution. Low inventory on 12 SKUs requires attention.",
        'customer_insights': "New York and California continue to be top markets. Average purchase latency of 15 days suggests effective marketing.",
        'financial_insights': "Debit card remains the preferred payment method. 5 unsettled orders require follow-up. Cash flow is stable.",
        'recommendations': [
            "Increase inventory for top-selling SKUs showing low stock",
            "Expand Amazon channel presence given strong performance",
            "Target marketing in NY and CA markets",
            "Review commission rates with eBay to improve net margins"
        ],
        'risk_alerts': [
            "12 SKUs below minimum stock threshold",
            "5 orders pending payment settlement"
        ]
    }
    
    test_data = {
        'report_month': 'October 2025',
        'sales_summary': {
            'total_gross': 224674.36,
            'total_orders': 350
        }
    }
    
    generator = ExecutiveSummaryGenerator()
    generator.generate_presentation(test_insights, test_data, 'test_executive_summary.pptx')
    print("Test PowerPoint generated successfully!")
