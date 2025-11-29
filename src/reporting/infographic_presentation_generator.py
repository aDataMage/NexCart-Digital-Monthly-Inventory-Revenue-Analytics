"""
Infographic-based PowerPoint generator for executive summaries.
Uses AI-generated infographics instead of text-heavy slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import logging
from .infographic_generator import InfographicGenerator

logger = logging.getLogger(__name__)


class InfographicPresentationGenerator:
    """Generate PowerPoint presentations with AI-generated infographics."""

    def __init__(self):
        """Initialize presentation generator."""
        self.prs = Presentation()
        # 16:9 aspect ratio for modern presentations
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        self.infographic_gen = InfographicGenerator()

        # Color scheme
        self.primary_color = RGBColor(52, 152, 219)
        self.text_color = RGBColor(44, 62, 80)

    def generate_presentation(
        self,
        insights: dict,
        report_data: dict,
        output_path: str,
        infographic_dir: str = None
    ):
        """
        Generate PowerPoint with infographic slides.

        Args:
            insights: LLM-generated insights dictionary
            report_data: Aggregated report data
            output_path: Path to save the .pptx file
            infographic_dir: Directory to save/load infographics
        """
        try:
            output_file = Path(output_path)
            if infographic_dir is None:
                infographic_dir = str(output_file.parent / 'infographics')

            logger.info("Generating infographics...")
            infographics = self.infographic_gen.generate_all_infographics(
                insights, report_data, infographic_dir
            )

            report_month = report_data.get('report_month', 'Monthly Report')

            # Create slides with infographics
            slide_configs = [
                ('title', f"Executive Summary - {report_month}"),
                ('executive_summary', "Executive Summary"),
                ('sales', "Sales & Channel Performance"),
                ('products', "Product & Inventory Performance"),
                ('customers', "Customer Demographics"),
                ('financial', "Financial Analysis"),
                ('recommendations', "Strategic Recommendations"),
                ('risks', "Risk Alerts"),
            ]

            for key, title in slide_configs:
                if key in infographics:
                    self._create_infographic_slide(title, infographics[key])
                else:
                    # Fallback to text slide if infographic failed
                    self._create_fallback_slide(title, insights, key)

            # Save presentation
            self.prs.save(output_path)
            logger.info(f"Infographic presentation saved to {output_path}")

            return infographics

        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            raise

    def _create_infographic_slide(self, title: str, image_path: str):
        """Create a slide with an infographic image."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Add infographic image (full slide)
        try:
            slide.shapes.add_picture(
                image_path,
                Inches(0), Inches(0),
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )
            logger.info(f"Added infographic slide: {title}")
        except Exception as e:
            logger.warning(f"Could not add image {image_path}: {e}")
            # Add title as fallback
            self._add_title_text(slide, title)

    def _create_fallback_slide(self, title: str, insights: dict, key: str):
        """Create a text-based fallback slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        background = slide.shapes.add_shape(
            1, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(245, 247, 250)
        background.line.fill.background()

        # Title
        self._add_title_text(slide, title)

        # Content based on key
        content = self._get_fallback_content(insights, key)
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5),
                Inches(12.333), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.text = content
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.font.size = Pt(16)
                para.font.color.rgb = self.text_color

    def _add_title_text(self, slide, title: str):
        """Add title text to a slide."""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(1)
        )
        tf = title_box.text_frame
        tf.text = title
        para = tf.paragraphs[0]
        para.font.size = Pt(32)
        para.font.bold = True
        para.font.color.rgb = self.primary_color

    def _get_fallback_content(self, insights: dict, key: str) -> str:
        """Get fallback text content for a slide."""
        mapping = {
            'executive_summary': insights.get('executive_summary', ''),
            'sales': insights.get('sales_insights', ''),
            'products': insights.get('product_insights', ''),
            'customers': insights.get('customer_insights', ''),
            'financial': insights.get('financial_insights', ''),
            'recommendations': '\n'.join([
                f"• {r}" for r in insights.get('recommendations', [])
            ]),
            'risks': '\n'.join([
                f"⚠ {r}" for r in insights.get('risk_alerts', [])
            ]),
        }
        return mapping.get(key, '')


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)

    test_insights = {
        'executive_summary': "Strong performance with 15% growth.",
        'sales_insights': "Amazon leads with 40% revenue share.",
        'product_insights': "Top 20% SKUs drive 80% of sales.",
        'customer_insights': "NY and CA are top markets.",
        'financial_insights': "95% settlement rate achieved.",
        'recommendations': ["Expand Amazon", "Restock top SKUs"],
        'risk_alerts': ["12 SKUs low stock", "5 pending orders"]
    }

    test_data = {
        'report_month': 'October 2025',
        'sales_summary': {'total_gross': 224674.36, 'total_orders': 350},
    }

    generator = InfographicPresentationGenerator()
    generator.generate_presentation(
        test_insights,
        test_data,
        'reports/test_infographic_presentation.pptx'
    )
    print("Test presentation generated!")
